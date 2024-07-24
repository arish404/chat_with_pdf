import streamlit as st
import PyPDF2
from re import sub
from docx import Document
from TexSoup import TexSoup
from pptx import Presentation

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text.strip()

def extract_text_from_pdf(uploaded_pdf):
    text = ""
    reader = PyPDF2.PdfReader(uploaded_pdf)

    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text += page.extract_text()

    return text

def extract_text_from_txt(uploaded_txt):
    file_contents = uploaded_txt.read()
    text = file_contents.decode('utf-8')
    return text

def extract_text_from_docx(file):
    doc = Document(file)
    text = ""

    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    return text

def extract_text_from_tex(tex_content):
    soup = TexSoup(tex_content)
    text = ""

    for node in soup.contents:
        if isinstance(node.text, list):
            text += ' '.join(node.text).strip() + "\n"

        else:
            text += node.text.strip() + "\n"

    return text.strip()

def preprocess_text(text):
    text = sub(r'\s+', ' ', text)
    text = text.lower()
    return text

def main():
    hide_streamlit_style = """
            <style>
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
            </style>
            """
    st.markdown(hide_streamlit_style, unsafe_allow_html=True)
    st.title('Text Extractor')
    uploaded_files = st.file_uploader("Upload a  file", accept_multiple_files=True)
    text =clean_text=" "

    if  uploaded_files:

        with st.form(key='query_form'):
            user_input = st.text_input(label=" ", placeholder="Enter prompts about your document", )
            submit_button = st.form_submit_button(label='Submit')

        if submit_button:
            for uploaded_file in uploaded_files:
                name = uploaded_file.name
                if( name.endswith(".pdf")):
                    text = extract_text_from_pdf( uploaded_file)

                elif(name.endswith(".docx") or name.endswith(".doc")):
                    text = extract_text_from_docx( uploaded_file)

                elif( name.endswith(".tex")):
                    tex_content = uploaded_file.read().decode('utf-8')
                    text = extract_text_from_tex( tex_content)

                elif( name.endswith(".pptx")):
                    text = extract_text_from_pptx( uploaded_file)
                
                elif( name.endswith(".txt")):
                    text = extract_text_from_txt( uploaded_file)

                clean_text += name+":"+preprocess_text(text)+"\n\n\n"
            st.header('Extracted Info')
            if user_input is not None:
                return clean_text,user_input 
            else:
                return clean_text,"  "
    else:
        st.write("Upload Your Documents")

if __name__ == "__main__":
    main()
